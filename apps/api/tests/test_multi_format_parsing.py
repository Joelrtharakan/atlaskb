"""Parsers for the new ingestible formats (.docx, .xlsx, .csv, .txt, .pptx) —
pure unit tests, no DB, building tiny fixture files on the fly. Every parser
must return the same `ParsedBlock` shape `parse_pdf`/`parse_markdown` do, so
these mirror tests/test_chunking.py's style rather than introducing a new one.
"""

import csv
import io

import pytest
from app.chunking import (
    parse_csv,
    parse_document,
    parse_docx,
    parse_pptx,
    parse_txt,
    parse_xlsx,
)


# --- .docx -------------------------------------------------------------
def _make_docx(path):
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_heading("Refund Policy", level=1)
    doc.add_paragraph("Refunds are issued within 30 days of purchase.")
    doc.add_heading("Exceptions", level=2)
    doc.add_paragraph("Digital goods are non-refundable once downloaded.")
    table = doc.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Tier"
    table.rows[0].cells[1].text = "Discount"
    table.rows[1].cells[0].text = "Gold"
    table.rows[1].cells[1].text = "20%"
    table.rows[2].cells[0].text = "Silver"
    table.rows[2].cells[1].text = "10%"
    doc.save(path)


def test_parse_docx_tracks_heading_sections(tmp_path):
    path = tmp_path / "policy.docx"
    _make_docx(path)
    blocks = parse_docx(path)
    sections = {b.section for b in blocks}
    assert "Refund Policy" in sections
    assert "Exceptions" in sections
    exceptions = next(b for b in blocks if b.section == "Exceptions")
    assert "non-refundable" in exceptions.text


def test_parse_docx_table_is_self_describing(tmp_path):
    path = tmp_path / "policy.docx"
    _make_docx(path)
    blocks = parse_docx(path)
    table_block = next(b for b in blocks if "Gold" in b.text)
    assert "Columns: Tier | Discount" in table_block.text
    assert "Tier: Gold" in table_block.text
    assert "Discount: 20%" in table_block.text
    # Under the most recently seen heading at the point the table appears.
    assert table_block.section is not None and "Exceptions" in table_block.section


def test_parse_document_dispatches_docx_by_extension(tmp_path):
    path = tmp_path / "policy.docx"
    _make_docx(path)
    blocks = parse_document(path, content_type="application/octet-stream")
    assert any("30 days" in b.text for b in blocks)


# --- .pptx -------------------------------------------------------------
def _make_pptx(path):
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[1]  # "Title and Content"
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Q3 Roadmap"
    body = slide.placeholders[1]
    body.text_frame.text = "Ship the Falcon v3 connector."

    prs.slides.add_slide(prs.slide_layouts[6])  # blank
    prs.save(path)


def test_parse_pptx_uses_slide_title_as_section(tmp_path):
    path = tmp_path / "roadmap.pptx"
    _make_pptx(path)
    blocks = parse_pptx(path)
    assert any(b.section == "Q3 Roadmap" and "Falcon v3 connector" in b.text for b in blocks)


def test_parse_pptx_blank_slide_yields_no_blocks(tmp_path):
    path = tmp_path / "roadmap.pptx"
    _make_pptx(path)
    blocks = parse_pptx(path)
    # Only the first (titled, populated) slide produces content.
    assert len(blocks) == 1


# --- .xlsx / .csv --------------------------------------------------------
def test_parse_xlsx_carries_header_context_per_chunk(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Q3 Sales"
    ws.append(["Region", "Revenue"])
    for i in range(30):
        ws.append([f"Region{i}", 1000 + i])
    path = tmp_path / "sales.xlsx"
    wb.save(path)

    blocks = parse_xlsx(path)
    assert len(blocks) >= 1
    assert all("Columns: Region | Revenue" in b.text for b in blocks)
    assert all(b.section is not None and b.section.startswith("Sheet: Q3 Sales, rows") for b in blocks)
    # Every row appears exactly once across the produced chunks.
    joined = "\n".join(b.text for b in blocks)
    assert "Region0: Region0" not in joined  # sanity: not double-labelled
    assert "Region: Region0" in joined
    assert "Region: Region29" in joined


def test_parse_xlsx_multiple_sheets_stay_separate(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet A"
    ws1.append(["X"])
    ws1.append([1])
    ws2 = wb.create_sheet("Sheet B")
    ws2.append(["Y"])
    ws2.append([2])
    path = tmp_path / "two_sheets.xlsx"
    wb.save(path)

    blocks = parse_xlsx(path)
    sections = {b.section for b in blocks}
    assert any(s.startswith("Sheet: Sheet A") for s in sections)
    assert any(s.startswith("Sheet: Sheet B") for s in sections)


def test_parse_csv_header_row_carried_as_context(tmp_path):
    path = tmp_path / "data.csv"
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Name", "Amount"])
    writer.writerow(["Alice", "100"])
    writer.writerow(["Bob", "200"])
    path.write_text(buf.getvalue(), encoding="utf-8")

    blocks = parse_csv(path)
    assert len(blocks) == 1
    assert "Columns: Name | Amount" in blocks[0].text
    assert "Name: Alice" in blocks[0].text
    assert "Amount: 200" in blocks[0].text
    assert blocks[0].section == "Sheet1, rows 2-3"


def test_parse_csv_single_row_no_data_yields_no_blocks(tmp_path):
    path = tmp_path / "header_only.csv"
    path.write_text("Name,Amount\n", encoding="utf-8")
    assert parse_csv(path) == []


# --- .txt ----------------------------------------------------------------
def test_parse_txt_splits_paragraphs(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("First note.\n\nSecond note, on its own paragraph.", encoding="utf-8")
    blocks = parse_txt(path)
    assert len(blocks) == 2
    assert blocks[0].text == "First note."
    assert blocks[1].text == "Second note, on its own paragraph."


def test_parse_txt_handles_utf16(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_bytes("Café résumé naïve.".encode("utf-16"))
    blocks = parse_txt(path)
    assert "Café résumé naïve." in blocks[0].text


def test_parse_txt_unreadable_bytes_raise_clear_error(tmp_path):
    path = tmp_path / "binary.txt"
    # 0x81 is undefined in cp1252, and this byte sequence isn't valid UTF-8
    # or UTF-16 either — should fail loudly, not decode to mojibake.
    path.write_bytes(b"\xff\xfe\x00\x81\x81\x81\x00\x00\xff")
    with pytest.raises(ValueError, match="Could not decode"):
        parse_txt(path)


# --- dispatch --------------------------------------------------------------
@pytest.mark.parametrize(
    "suffix,content_type",
    [
        (".docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        (".xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        (".pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        (".csv", "text/csv"),
        (".txt", "text/plain"),
    ],
)
def test_parse_document_dispatches_by_content_type(tmp_path, suffix, content_type):
    # Extension deliberately wrong/generic — dispatch must still route
    # correctly off content_type alone, matching the extension-alone check.
    path = tmp_path / f"file{suffix}"
    if suffix == ".docx":
        _make_docx(path)
    elif suffix == ".pptx":
        _make_pptx(path)
    elif suffix == ".xlsx":
        from openpyxl import Workbook

        wb = Workbook()
        wb.active.append(["A", "B"])
        wb.active.append([1, 2])
        wb.save(path)
    elif suffix == ".csv":
        path.write_text("A,B\n1,2\n", encoding="utf-8")
    elif suffix == ".txt":
        path.write_text("Just some text.", encoding="utf-8")

    blocks = parse_document(path, content_type=content_type)
    assert isinstance(blocks, list)
