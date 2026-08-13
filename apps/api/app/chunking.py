"""Document parsing and chunking.

Everything here is pure and side-effect free (aside from reading the source file
in ``parse_document``) so the chunking logic can be unit-tested without a
database, a worker, or an embedding model.

Parsing produces a list of :class:`ParsedBlock` — a span of text tagged with the
location metadata we want to preserve for citations (page number for PDFs,
section heading for Markdown/HTML). Chunking then packs those blocks into
overlapping windows, carrying the metadata through onto each chunk.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class ParsedBlock:
    text: str
    page_num: int | None = None
    section: str | None = None


@dataclass(frozen=True)
class ChunkData:
    chunk_index: int
    text: str
    page_num: int | None = None
    section: str | None = None


# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------
_BLANK_LINE_RE = re.compile(r"\n\s*\n+")


def _split_paragraphs(text: str) -> list[str]:
    """Split a page's extracted text into paragraphs on blank lines.

    PDF text-extraction quality varies: some PDFs preserve blank lines between
    paragraphs, others collapse a page to one run of lines with no blank-line
    boundaries at all. When no boundary exists, the page stays one paragraph —
    ``chunk_blocks``' word-window packing still applies to it exactly as before,
    so this is additive: it only stops splitting *across* paragraph boundaries
    when the source actually has them, instead of always word-wrapping blind.
    """
    paras = [p.strip() for p in _BLANK_LINE_RE.split(text) if p.strip()]
    return paras or ([text] if text.strip() else [])


def parse_pdf(path: str | Path) -> list[ParsedBlock]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    blocks: list[ParsedBlock] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        for para in _split_paragraphs(text):
            blocks.append(ParsedBlock(text=para, page_num=i))
    return blocks


_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")


def parse_markdown(text: str) -> list[ParsedBlock]:
    """Split Markdown into paragraph blocks, tagging each with its section.

    The section is the nearest preceding ATX heading (``#`` .. ``######``).
    """
    blocks: list[ParsedBlock] = []
    section: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        para = "\n".join(buffer).strip()
        if para:
            blocks.append(ParsedBlock(text=para, section=section))
        buffer = []

    for line in text.splitlines():
        heading = _HEADING_RE.match(line.strip())
        if heading:
            flush()
            section = heading.group(2).strip()
        elif line.strip() == "":
            flush()
        else:
            buffer.append(line)
    flush()
    return blocks


def parse_html(html: str) -> list[ParsedBlock]:
    """Extract visible text from HTML, tagging blocks with their section.

    Section tracking follows the most recent ``<h1>``..``<h6>`` heading.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    body = soup.body or soup
    blocks: list[ParsedBlock] = []
    section: str | None = None

    block_tags = ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "blockquote"]
    elements = body.find_all(block_tags)
    if not elements:
        text = soup.get_text(separator="\n").strip()
        return [ParsedBlock(text=text)] if text else []

    for el in elements:
        text = el.get_text(separator=" ", strip=True)
        if not text:
            continue
        if el.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            section = text
        else:
            blocks.append(ParsedBlock(text=text, section=section))
    return blocks


def _read_text_smart(path: str | Path) -> str:
    """Decode a text file without assuming UTF-8. Tries UTF-8, then UTF-16
    (common from Windows/Excel "Save as Unicode Text"), then CP1252 (common
    from older Windows tools) — all strict decodes, so a real mismatch raises
    rather than silently producing garbled/mojibake text. Deliberately does
    *not* fall back to Latin-1: it maps every byte value to *something*, so it
    never raises and would defeat the point of this function — a clear error
    beats a document full of silently-wrong text."""
    raw = Path(path).read_bytes()
    for encoding in ("utf-8", "utf-16", "cp1252"):
        try:
            return raw.decode(encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(
        f"Could not decode {Path(path).name} as UTF-8, UTF-16, or CP1252 text."
    )


def parse_txt(path: str | Path) -> list[ParsedBlock]:
    text = _read_text_smart(path)
    return [ParsedBlock(text=p) for p in _split_paragraphs(text)]


def parse_docx(path: str | Path) -> list[ParsedBlock]:
    """Heading styles ("Heading 1"/"Heading 2"/... or "Title") become the
    running `section` for the paragraphs under them, same as Markdown's ATX
    headings and HTML's h1-h6. Tables are walked in document order (not
    collected separately at the end) and rendered through the same
    self-describing row-range shape spreadsheet sheets use, so a table chunk
    never loses its column headers."""
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = DocxDocument(str(path))
    blocks: list[ParsedBlock] = []
    section: str | None = None

    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or "") if para.style is not None else ""
            if style_name.lower().startswith("heading") or style_name.lower() == "title":
                section = text
            else:
                blocks.append(ParsedBlock(text=text, section=section))
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            blocks.extend(_chunk_rows(rows, section=section, label=None))

    return blocks


def parse_pptx(path: str | Path) -> list[ParsedBlock]:
    """Each slide is a section (its title, or "Slide N" if titleless); body
    text boxes become blocks under it. Embedded images/charts are skipped —
    text content only, a stated boundary (see the multi-format ingestion
    scope note), not a silent gap."""
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[ParsedBlock] = []

    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = title_shape.text_frame.text.strip() if title_shape is not None else ""
        section = title_text or f"Slide {i}"

        title_id = title_shape.shape_id if title_shape is not None else None
        body_found = False
        for shape in slide.shapes:
            # `shape is title_shape` doesn't reliably hold — python-pptx
            # wraps the same underlying XML element in a fresh object on
            # each access, so identity comparison is unreliable; compare
            # the stable shape_id instead.
            if shape.shape_id == title_id or not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            body_found = True
            blocks.append(ParsedBlock(text=text, section=section))
        if not body_found and title_text:
            blocks.append(ParsedBlock(text=title_text, section=section))

    return blocks


def _chunk_rows(
    rows: list[list[str]],
    *,
    section: str | None,
    label: str | None,
    max_chars: int = 800,
) -> list[ParsedBlock]:
    """Row-range blocks with the header row's column names rendered into
    every block's own text (as "Columns: ..." plus "col: value" per row) —
    not just attached as separate metadata — so a chunk retrieved on its own
    is self-describing rather than a bag of values with no column names.
    Splits by character budget rather than a fixed row count, so very wide
    rows still produce reasonably-sized chunks. ``label`` is prefixed to the
    row-range section string (e.g. "Sheet: Q3 Sales") — a table with no
    sheet name (a Word table) omits it and reads as "rows 2-14"."""
    non_empty = [r for r in rows if any(str(v).strip() for v in r)]
    if len(non_empty) < 2:  # need a header row plus at least one data row
        return []

    header = [str(v).strip() for v in non_empty[0]]
    header_line = "Columns: " + " | ".join(header)
    data_rows = non_empty[1:]

    def row_range_section(start: int, end: int) -> str:
        prefix = f"{label}, " if label else ""
        combined = f"{prefix}rows {start}-{end}"
        return f"{section} — {combined}" if section else combined

    blocks: list[ParsedBlock] = []
    lines: list[str] = []
    range_start = 2  # 1-based; row 1 is the header
    length = len(header_line)

    def flush(range_end: int) -> None:
        if not lines:
            return
        text = header_line + "\n" + "\n".join(lines)
        blocks.append(ParsedBlock(text=text, section=row_range_section(range_start, range_end)))

    for offset, row in enumerate(data_rows):
        row_num = offset + 2
        cells = [str(v).strip() for v in row]
        line = " | ".join(f"{h}: {v}" for h, v in zip(header, cells))
        if lines and length + len(line) + 1 > max_chars:
            flush(row_num - 1)
            lines = []
            range_start = row_num
            length = len(header_line)
        lines.append(line)
        length += len(line) + 1

    flush(range_start + len(lines) - 1)
    return blocks


def parse_xlsx(path: str | Path) -> list[ParsedBlock]:
    """Each sheet is its own logical section — a chunk never blends rows from
    two different sheets, so "rows 2-14" is unambiguous without also naming
    the sheet in every single row."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        blocks: list[ParsedBlock] = []
        for ws in wb.worksheets:
            rows = [
                ["" if v is None else v for v in row]
                for row in ws.iter_rows(values_only=True)
            ]
            blocks.extend(_chunk_rows(rows, section=None, label=f"Sheet: {ws.title}"))
        return blocks
    finally:
        wb.close()


def parse_csv(path: str | Path) -> list[ParsedBlock]:
    """Treated as .xlsx's single-sheet case: the whole file is one sheet, the
    first row is the header, and chunks carry it the same self-describing way."""
    import csv
    import io

    text = _read_text_smart(path)
    rows = list(csv.reader(io.StringIO(text)))
    return _chunk_rows(rows, section=None, label="Sheet1")


def parse_document(path: str | Path, content_type: str) -> list[ParsedBlock]:
    """Dispatch to the right parser based on content type / file extension —
    every parser above returns the same `ParsedBlock` shape, so nothing
    downstream of this function (chunking, embedding, write, retrieval) forks
    per format."""
    path = Path(path)
    ct = (content_type or "").lower()
    suffix = path.suffix.lower()

    if "pdf" in ct or suffix == ".pdf":
        return parse_pdf(path)
    if "wordprocessingml.document" in ct or suffix == ".docx":
        return parse_docx(path)
    if "presentationml.presentation" in ct or suffix == ".pptx":
        return parse_pptx(path)
    if "spreadsheetml.sheet" in ct or suffix == ".xlsx":
        return parse_xlsx(path)
    if ct == "text/csv" or suffix == ".csv":
        return parse_csv(path)
    if suffix == ".txt" or ct == "text/plain":
        return parse_txt(path)

    raw = path.read_text(encoding="utf-8", errors="replace")
    if "html" in ct or suffix in (".html", ".htm"):
        return parse_html(raw)
    if "markdown" in ct or ct == "text/md" or suffix in (".md", ".markdown"):
        return parse_markdown(raw)
    # Fallback: treat as Markdown/plain text.
    return parse_markdown(raw)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------
_WORD_RE = re.compile(r"\S+")


def _split_with_overlap(text: str, max_chars: int, overlap: int) -> list[str]:
    """Pack whitespace-delimited tokens into windows of at most ``max_chars``.

    Consecutive windows share roughly ``overlap`` characters of trailing context
    so a fact spanning a chunk boundary is still retrievable from both sides.
    """
    words = _WORD_RE.findall(text)
    if not words:
        return []

    chunks: list[str] = []
    current: list[str] = []
    length = 0  # characters in ``current`` including single-space separators

    for word in words:
        add = len(word) + (1 if current else 0)
        if current and length + add > max_chars:
            chunks.append(" ".join(current))
            # Seed the next window with trailing words up to ``overlap`` chars.
            carry: list[str] = []
            carry_len = 0
            for w in reversed(current):
                extra = len(w) + (1 if carry else 0)
                if carry_len + extra > overlap:
                    break
                carry.insert(0, w)
                carry_len += extra
            current = carry
            length = carry_len
            add = len(word) + (1 if current else 0)
        current.append(word)
        length += add

    if current:
        chunks.append(" ".join(current))
    return chunks


def chunk_blocks(
    blocks: list[ParsedBlock],
    *,
    max_chars: int = 1000,
    overlap: int = 150,
) -> list[ChunkData]:
    """Turn parsed blocks into ordered, overlapping chunks.

    Each block is chunked independently so page/section metadata is never mixed
    across boundaries. ``chunk_index`` is a stable global ordering.
    """
    if overlap >= max_chars:
        raise ValueError("overlap must be smaller than max_chars")

    out: list[ChunkData] = []
    index = 0
    for block in blocks:
        for piece in _split_with_overlap(block.text, max_chars, overlap):
            out.append(
                ChunkData(
                    chunk_index=index,
                    text=piece,
                    page_num=block.page_num,
                    section=block.section,
                )
            )
            index += 1
    return out
